import json
from io import BytesIO
from typing import cast

import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pypdf import PdfWriter

from harness.core.models import ProcessingStatus
from harness.inputs.processors import DefaultInputProcessor


def docx_bytes() -> bytes:
    document = Document()
    document.add_heading("Quarterly Report", level=1)
    document.add_paragraph("Revenue increased.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "42"
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Summary"
    sheet.append(["Metric", "Value"])
    sheet.append(["Revenue", 42])
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    assert title is not None and title.has_text_frame
    assert title.text_frame is not None
    title.text_frame.text = "Launch Plan"
    body = cast(SlidePlaceholder, slide.placeholders[1])
    assert body.has_text_frame and body.text_frame is not None
    body.text_frame.text = "Ship in July"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def pdf_bytes() -> bytes:
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def png_bytes() -> bytes:
    output = BytesIO()
    Image.new("RGB", (3, 2), color="red").save(output, format="PNG")
    return output.getvalue()


@pytest.mark.parametrize(
    ("name", "media_type", "content", "processor", "needle"),
    [
        ("notes.txt", "text/plain", b"hello\r\nworld", "text", "hello\nworld"),
        ("data.json", "application/json", b'{"b":2,"a":1}', "text", '"a": 1'),
        (
            "report.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            docx_bytes(),
            "docx",
            "# Quarterly Report",
        ),
        (
            "metrics.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            xlsx_bytes(),
            "xlsx",
            "Revenue",
        ),
        (
            "deck.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            pptx_bytes(),
            "pptx",
            "Launch Plan",
        ),
        ("paper.pdf", "application/pdf", pdf_bytes(), "pdf", "## Page 1"),
    ],
)
def test_documents_produce_deterministic_markdown_or_text_projection(
    name: str, media_type: str, content: bytes, processor: str, needle: str
) -> None:
    result = DefaultInputProcessor().process(
        name=name, media_type=media_type, content=content
    )

    assert result.status is ProcessingStatus.PROCESSED
    assert result.processor == processor
    assert result.derived
    assert needle in result.derived[0].content.decode("utf-8")
    assert all(".." not in file.relative_path for file in result.derived)


def test_image_is_validated_and_reports_dimensions_without_mutating_bytes() -> None:
    content = png_bytes()

    result = DefaultInputProcessor().process(
        name="pixel.png", media_type="image/png", content=content
    )

    assert result.status is ProcessingStatus.PROCESSED
    assert result.processor == "image"
    assert result.metadata == {"format": "PNG", "width": 3, "height": 2}
    assert result.derived == ()
    assert content == png_bytes()


def test_unsupported_input_is_explicit_and_keeps_original_available() -> None:
    result = DefaultInputProcessor().process(
        name="archive.bin", media_type="application/octet-stream", content=b"raw"
    )

    assert result.status is ProcessingStatus.UNSUPPORTED
    assert result.processor == "unsupported"
    assert result.derived == ()
    assert json.dumps(result.metadata) == "{}"
