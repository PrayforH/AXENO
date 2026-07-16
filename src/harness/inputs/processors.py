"""Deterministic, local processors for common browser-uploaded formats."""

import csv
import json
import re
from collections.abc import Mapping
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any, cast

from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.picture import Picture
from pypdf import PdfReader

from harness.core.models import ProcessingStatus
from harness.inputs.base import DerivedInput, InputProcessingResult


def _stem(name: str) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]", "_", Path(name).stem).strip("._")
    return value[:80] or "input"


def _normalized_text(content: bytes) -> str:
    return content.decode("utf-8-sig").replace("\r\n", "\n").replace("\r", "\n")


def _markdown_table(rows: list[list[str]]) -> list[str]:
    if not rows:
        return []
    width = max(len(row) for row in rows)
    padded = [row + [""] * (width - len(row)) for row in rows]
    escaped = [[cell.replace("|", "\\|").replace("\n", " ") for cell in row] for row in padded]
    return [
        "| " + " | ".join(escaped[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
        *("| " + " | ".join(row) + " |" for row in escaped[1:]),
    ]


class DefaultInputProcessor:
    def process(
        self, *, name: str, media_type: str, content: bytes
    ) -> InputProcessingResult:
        suffix = Path(name).suffix.lower()
        if media_type.startswith("image/") or suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return self._image(content)
        if suffix == ".docx" or "wordprocessingml" in media_type:
            return self._docx(name, content)
        if suffix == ".xlsx" or "spreadsheetml" in media_type:
            return self._xlsx(name, content)
        if suffix == ".pptx" or "presentationml" in media_type:
            return self._pptx(name, content)
        if suffix == ".pdf" or media_type == "application/pdf":
            return self._pdf(name, content)
        if (
            media_type.startswith("text/")
            or suffix in {".txt", ".md", ".csv", ".json", ".xml"}
            or media_type in {"application/json", "application/xml"}
        ):
            return self._text(name, media_type, content)
        return InputProcessingResult(
            status=ProcessingStatus.UNSUPPORTED,
            processor="unsupported",
        )

    def _text(self, name: str, media_type: str, content: bytes) -> InputProcessingResult:
        text = _normalized_text(content)
        suffix = Path(name).suffix.lower()
        if suffix == ".json" or media_type == "application/json":
            parsed: Any = json.loads(text)
            text = json.dumps(parsed, ensure_ascii=False, indent=2, sort_keys=True)
            media_type = "application/json"
            extension = "json"
        else:
            extension = "md" if suffix == ".md" else "txt"
            media_type = "text/markdown" if extension == "md" else "text/plain"
        return InputProcessingResult(
            status=ProcessingStatus.PROCESSED,
            processor="text",
            derived=(
                DerivedInput(
                    relative_path=f"{_stem(name)}.{extension}",
                    media_type=media_type,
                    content=text.encode("utf-8"),
                ),
            ),
        )

    def _docx(self, name: str, content: bytes) -> InputProcessingResult:
        document = Document(BytesIO(content))
        lines: list[str] = []
        headings: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style = (paragraph.style.name or "") if paragraph.style is not None else ""
            match = re.match(r"Heading (\d+)", style)
            if match:
                level = min(int(match.group(1)), 6)
                lines.append(f"{'#' * level} {text}")
                headings.append(text)
            else:
                lines.append(text)
        for table in document.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            lines.extend(["", *_markdown_table(rows)])
        derived: list[DerivedInput] = [
            DerivedInput(
                relative_path=f"{_stem(name)}.md",
                media_type="text/markdown",
                content=("\n\n".join(lines).strip() + "\n").encode("utf-8"),
                metadata={"headings": headings},
            )
        ]
        for index, relationship in enumerate(
            sorted(document.part.rels.values(), key=lambda item: item.rId), start=1
        ):
            if relationship.is_external:
                continue
            target = relationship.target_part
            if not getattr(target, "content_type", "").startswith("image/"):
                continue
            filename = Path(relationship.target_ref).name
            derived.append(
                DerivedInput(
                    relative_path=f"images/{index:02d}-{filename}",
                    media_type=target.content_type,
                    content=target.blob,
                )
            )
        return InputProcessingResult(
            status=ProcessingStatus.PROCESSED,
            processor="docx",
            derived=tuple(derived),
            metadata={"headings": headings, "tables": len(document.tables)},
        )

    def _xlsx(self, name: str, content: bytes) -> InputProcessingResult:
        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        lines = [f"# Workbook: {Path(name).name}"]
        derived: list[DerivedInput] = []
        for worksheet in workbook.worksheets:
            rows = [
                ["" if value is None else str(value) for value in row]
                for row in worksheet.iter_rows(values_only=True)
            ]
            while rows and not any(rows[-1]):
                rows.pop()
            lines.extend(["", f"## {worksheet.title}", *_markdown_table(rows)])
            output = StringIO(newline="")
            writer = csv.writer(output, lineterminator="\n")
            writer.writerows(rows)
            derived.append(
                DerivedInput(
                    relative_path=f"sheets/{_stem(worksheet.title)}.csv",
                    media_type="text/csv",
                    content=output.getvalue().encode("utf-8"),
                    metadata={"sheet": worksheet.title},
                )
            )
        summary = DerivedInput(
            relative_path=f"{_stem(name)}.md",
            media_type="text/markdown",
            content=("\n".join(lines).strip() + "\n").encode("utf-8"),
        )
        return InputProcessingResult(
            status=ProcessingStatus.PROCESSED,
            processor="xlsx",
            derived=(summary, *derived),
            metadata={"sheets": workbook.sheetnames},
        )

    def _pptx(self, name: str, content: bytes) -> InputProcessingResult:
        presentation = Presentation(BytesIO(content))
        lines = [f"# Presentation: {Path(name).name}"]
        derived: list[DerivedInput] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines.extend(["", f"## Slide {slide_index}"])
            for shape_index, shape in enumerate(slide.shapes, start=1):
                text = getattr(shape, "text", "").strip()
                if text:
                    lines.append(text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = cast(Picture, shape).image
                    derived.append(
                        DerivedInput(
                            relative_path=(
                                f"images/slide-{slide_index:02d}-{shape_index:02d}."
                                f"{image.ext}"
                            ),
                            media_type=image.content_type,
                            content=image.blob,
                        )
                    )
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes = notes_frame.text.strip() if notes_frame is not None else ""
                if notes:
                    lines.extend(["", "### Speaker notes", notes])
        summary = DerivedInput(
            relative_path=f"{_stem(name)}.md",
            media_type="text/markdown",
            content=("\n".join(lines).strip() + "\n").encode("utf-8"),
        )
        return InputProcessingResult(
            status=ProcessingStatus.PROCESSED,
            processor="pptx",
            derived=(summary, *derived),
            metadata={"slides": len(presentation.slides)},
        )

    def _pdf(self, name: str, content: bytes) -> InputProcessingResult:
        reader = PdfReader(BytesIO(content))
        lines = [f"# PDF: {Path(name).name}"]
        for index, page in enumerate(reader.pages, start=1):
            lines.extend(["", f"## Page {index}", page.extract_text() or ""])
        raw_metadata = cast(Mapping[str, object], reader.metadata or {})
        metadata = {
            str(key).lstrip("/"): str(value)
            for key, value in raw_metadata.items()
        }
        return InputProcessingResult(
            status=ProcessingStatus.PROCESSED,
            processor="pdf",
            derived=(
                DerivedInput(
                    relative_path=f"{_stem(name)}.md",
                    media_type="text/markdown",
                    content=("\n".join(lines).strip() + "\n").encode("utf-8"),
                    metadata={"pages": len(reader.pages)},
                ),
            ),
            metadata={"pages": len(reader.pages), **metadata},
        )

    def _image(self, content: bytes) -> InputProcessingResult:
        with Image.open(BytesIO(content)) as image:
            image.verify()
        with Image.open(BytesIO(content)) as image:
            metadata = {
                "format": image.format or "unknown",
                "width": image.width,
                "height": image.height,
            }
        return InputProcessingResult(
            status=ProcessingStatus.PROCESSED,
            processor="image",
            metadata=metadata,
        )
